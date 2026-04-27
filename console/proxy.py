#!/usr/bin/env python3
"""
APEX Console proxy server.
  /api/messages          -> proxy to Anthropic API
  /api/generate-docx     -> markdown -> branded Word document
  /api/generate-pptx     -> markdown -> branded PowerPoint deck
  /api/engagements       -> SQLite CRUD for engagements
  /api/deliverables      -> SQLite CRUD for deliverable library
"""

import http.server
import urllib.request
import urllib.error
import urllib.parse
import json
import io
import os
import re
import sys
import sqlite3
import datetime

PORT          = 3000
ANTHROPIC_URL = "https://api.anthropic.com/v1/messages"
STATIC_DIR    = os.path.dirname(os.path.abspath(__file__))
DB_PATH       = os.path.join(STATIC_DIR, 'apex.db')

# APEX brand palette (R, G, B)
NAVY  = (15,  27,  45)
TEAL  = (14,  124, 123)
GOLD  = (201, 168, 76)
WHITE = (240, 245, 250)
FOG   = (139, 163, 190)
SLATE = (30,  45,  64)


# ── Database ──────────────────────────────────────────────────────────────────

def init_db():
    con = sqlite3.connect(DB_PATH)
    cur = con.cursor()
    cur.executescript("""
        CREATE TABLE IF NOT EXISTS engagements (
            id               TEXT PRIMARY KEY,
            client           TEXT NOT NULL,
            sector           TEXT,
            deliverable      TEXT,
            status           TEXT DEFAULT 'ready',
            brief            TEXT,
            agents           TEXT,
            created_at       TEXT DEFAULT (datetime('now')),
            updated_at       TEXT DEFAULT (datetime('now'))
        );

        CREATE TABLE IF NOT EXISTS deliverable_library (
            id               TEXT PRIMARY KEY,
            engagement_id    TEXT,
            gate_id          TEXT,
            label            TEXT NOT NULL,
            content          TEXT NOT NULL,
            verdict          TEXT,
            client           TEXT,
            sector           TEXT,
            deliverable_type TEXT,
            word_count       INTEGER DEFAULT 0,
            created_at       TEXT DEFAULT (datetime('now'))
        );

        CREATE INDEX IF NOT EXISTS idx_del_engagement
            ON deliverable_library(engagement_id);
        CREATE INDEX IF NOT EXISTS idx_del_verdict
            ON deliverable_library(verdict);
        CREATE INDEX IF NOT EXISTS idx_del_client
            ON deliverable_library(client);
    """)
    con.commit()
    con.close()
    print(f"  Database -> {DB_PATH}")


def db_connect():
    con = sqlite3.connect(DB_PATH)
    con.row_factory = sqlite3.Row
    return con


def row_to_dict(row):
    return dict(row) if row else None


# ── Markdown helpers ──────────────────────────────────────────────────────────

def _slug(text):
    return re.sub(r'[^a-z0-9]+', '-', text.lower()).strip('-')


def _strip_inline(text):
    """Remove **/*/`/[] markdown, leave plain text."""
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'\*(.+?)\*',     r'\1', text)
    text = re.sub(r'\[(.+?)\]\(.+?\)', r'\1', text)
    text = re.sub(r'`(.+?)`',       r'\1', text)
    return text


def _parse_md(content):
    """Return list of (type, text) tuples from markdown string."""
    blocks = []
    for raw in content.split('\n'):
        line = raw.rstrip()
        if   line.startswith('# '):    blocks.append(('title',  line[2:].strip()))
        elif line.startswith('## '):   blocks.append(('h1',     line[3:].strip()))
        elif line.startswith('### '):  blocks.append(('h2',     line[4:].strip()))
        elif line.startswith(('- ', '* ')): blocks.append(('bullet', line[2:].strip()))
        elif line == '---':            blocks.append(('hr',     ''))
        elif line == '':               blocks.append(('empty',  ''))
        else:                          blocks.append(('para',   line))
    return blocks


def _is_meta(text):
    return bool(re.match(r'\*?\*?Client:', text) or
                re.match(r'\*?\*?Sector:', text) or
                re.match(r'Date:', text))


# ── Word (.docx) builder ──────────────────────────────────────────────────────

def _build_docx(content, title, client, sector):
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor

    doc = Document()

    for sec in doc.sections:
        sec.top_margin    = Inches(1.0)
        sec.bottom_margin = Inches(1.0)
        sec.left_margin   = Inches(1.25)
        sec.right_margin  = Inches(1.25)

    def rgb(*t):
        return RGBColor(*t)

    def run(para, text, bold=False, color=None, size=None, italic=False):
        r = para.add_run(text)
        r.bold   = bold
        r.italic = italic
        if color: r.font.color.rgb = rgb(*color)
        if size:  r.font.size = Pt(size)
        return r

    def inline(para, text, size=11, default_color=None):
        parts = re.split(r'\*\*(.+?)\*\*', text)
        for i, part in enumerate(parts):
            if part:
                run(para, part, bold=(i % 2 == 1), size=size, color=default_color)

    cover = doc.add_paragraph()
    cover.alignment = 1
    run(cover, title, bold=True, color=TEAL, size=28)

    sub = doc.add_paragraph()
    sub.alignment = 1
    meta_text = '  ·  '.join(filter(None, [client, sector]))
    run(sub, meta_text, color=FOG, size=12)
    doc.add_paragraph()

    for btype, btext in _parse_md(content):
        if btype == 'title':
            pass
        elif btype == 'h1':
            h = doc.add_heading(level=1)
            h.clear()
            run(h, _strip_inline(btext), bold=True, color=TEAL, size=16)
        elif btype == 'h2':
            h = doc.add_heading(level=2)
            h.clear()
            run(h, _strip_inline(btext), bold=True, color=GOLD, size=13)
        elif btype == 'bullet':
            p = doc.add_paragraph(style='List Bullet')
            inline(p, _strip_inline(btext), size=11)
        elif btype == 'hr':
            doc.add_paragraph()
        elif btype in ('empty', ):
            pass
        elif btype == 'para':
            if _is_meta(btext):
                continue
            p = doc.add_paragraph()
            inline(p, _strip_inline(btext), size=11)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.read()


# ── PowerPoint (.pptx) builder ────────────────────────────────────────────────

def _build_pptx(content, title, client, sector):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def rgb(*t):
        return RGBColor(*t)

    def bg(slide, color=NAVY):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = rgb(*color)

    def accent_bar(slide):
        s = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.07), Inches(7.5))
        s.fill.solid()
        s.fill.fore_color.rgb = rgb(*TEAL)
        s.line.fill.background()

    def add_textbox(slide, l, t, w, h):
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tb.text_frame.word_wrap = True
        return tb.text_frame

    def para(tf, text, size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT, bullet=False):
        p   = tf.add_paragraph()
        p.alignment = align
        if bullet:
            p.level = 1
        r   = p.add_run()
        r.text = text
        r.font.size  = Pt(size)
        r.font.bold  = bold
        r.font.color.rgb = rgb(*color)
        return p

    sl = prs.slides.add_slide(blank)
    bg(sl); accent_bar(sl)

    tf_title = add_textbox(sl, 0.5, 2.2, 12.5, 1.6)
    para(tf_title, title, size=38, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

    tf_meta = add_textbox(sl, 0.5, 3.9, 12.5, 0.7)
    meta_text = '  ·  '.join(filter(None, [client, sector]))
    para(tf_meta, meta_text, size=16, color=FOG)

    tf_brand = add_textbox(sl, 0.5, 6.8, 4, 0.4)
    para(tf_brand, 'APEX  ·  AI-native strategy consulting', size=10, color=FOG)

    sections = []
    cur_h, cur_body = None, []
    for btype, btext in _parse_md(content):
        if btype == 'h1':
            if cur_h is not None:
                sections.append((cur_h, cur_body))
            cur_h, cur_body = btext, []
        elif btype in ('title', 'hr', 'empty'):
            continue
        elif cur_h is not None:
            cur_body.append((btype, btext))
    if cur_h is not None:
        sections.append((cur_h, cur_body))

    MAX_CHARS = 1400
    for heading, body_lines in sections:
        sl = prs.slides.add_slide(blank)
        bg(sl); accent_bar(sl)

        tf_h = add_textbox(sl, 0.35, 0.25, 12.7, 0.85)
        para(tf_h, _strip_inline(heading), size=22, bold=True, color=TEAL)

        bar = sl.shapes.add_shape(1, Inches(0.35), Inches(1.05), Inches(12.7), Pt(2))
        bar.fill.solid(); bar.fill.fore_color.rgb = rgb(*GOLD)
        bar.line.fill.background()

        tf_body = add_textbox(sl, 0.35, 1.25, 12.7, 5.9)
        chars = 0
        for btype, btext in body_lines:
            if chars > MAX_CHARS:
                break
            clean = _strip_inline(btext).strip()
            if not clean or _is_meta(clean):
                continue
            if btype == 'bullet':
                para(tf_body, '• ' + clean, size=14, color=WHITE, bullet=True)
            elif btype == 'h2':
                para(tf_body, clean, size=14, bold=True, color=GOLD)
            elif btype == 'para':
                para(tf_body, clean, size=13, color=WHITE)
            chars += len(clean)

    sl = prs.slides.add_slide(blank)
    bg(sl); accent_bar(sl)
    tf_end = add_textbox(sl, 0.5, 3.0, 12.5, 1.2)
    para(tf_end, 'APEX', size=48, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    tf_sub = add_textbox(sl, 0.5, 4.1, 12.5, 0.6)
    para(tf_sub, 'AI-native strategy consulting', size=16, color=FOG, align=PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ── HTTP handler ──────────────────────────────────────────────────────────────

class ProxyHandler(http.server.SimpleHTTPRequestHandler):

    def __init__(self, *args, **kwargs):
        super().__init__(*args, directory=STATIC_DIR, **kwargs)

    def log_message(self, format, *args):
        print(f"  {self.address_string()} {format % args}")

    def do_OPTIONS(self):
        self.send_response(200)
        self._cors_headers()
        self.end_headers()

    # ── GET ───────────────────────────────────────────────────────────────────

    def do_GET(self):
        if self.path == '/api/engagements':
            self._get_engagements()
        elif self.path.startswith('/api/deliverables'):
            self._get_deliverables()
        else:
            super().do_GET()

    # ── POST ──────────────────────────────────────────────────────────────────

    def do_POST(self):
        if   self.path == '/api/messages':       self._proxy_anthropic()
        elif self.path == '/api/generate-docx':  self._generate('docx')
        elif self.path == '/api/generate-pptx':  self._generate('pptx')
        elif self.path == '/api/engagements':    self._save_engagement()
        elif self.path == '/api/deliverables':   self._save_deliverable()
        else:                                     self.send_error(404)

    # ── DELETE ────────────────────────────────────────────────────────────────

    def do_DELETE(self):
        m = re.match(r'^/api/deliverables/(.+)$', self.path)
        if m:
            self._delete_deliverable(urllib.parse.unquote(m.group(1)))
        else:
            self.send_error(404)

    # ── CORS ──────────────────────────────────────────────────────────────────

    def _cors_headers(self):
        self.send_header('Access-Control-Allow-Origin',  '*')
        self.send_header('Access-Control-Allow-Methods', 'POST, GET, DELETE, OPTIONS')
        self.send_header('Access-Control-Allow-Headers',
                         'Content-Type, x-api-key, anthropic-version, anthropic-beta')

    def _json_response(self, data, status=200):
        body = json.dumps(data, default=str).encode()
        self.send_response(status)
        self._cors_headers()
        self.send_header('Content-Type',   'application/json')
        self.send_header('Content-Length', str(len(body)))
        self.end_headers()
        self.wfile.write(body)

    def _read_json_body(self):
        length = int(self.headers.get('Content-Length', 0))
        return json.loads(self.rfile.read(length))

    # ── Engagements API ───────────────────────────────────────────────────────

    def _get_engagements(self):
        try:
            con = db_connect()
            rows = con.execute(
                "SELECT * FROM engagements ORDER BY created_at DESC"
            ).fetchall()
            con.close()
            self._json_response([dict(r) for r in rows])
        except Exception as e:
            self._json_response({'error': str(e)}, 500)

    def _save_engagement(self):
        try:
            p = self._read_json_body()
            now = datetime.datetime.utcnow().isoformat()
            con = db_connect()
            con.execute("""
                INSERT INTO engagements
                    (id, client, sector, deliverable, status, brief, agents, created_at, updated_at)
                VALUES
                    (:id,:client,:sector,:deliverable,:status,:brief,:agents,:created_at,:updated_at)
                ON CONFLICT(id) DO UPDATE SET
                    status     = excluded.status,
                    updated_at = excluded.updated_at
            """, {
                'id':          p.get('id'),
                'client':      p.get('client', ''),
                'sector':      p.get('sector', ''),
                'deliverable': p.get('deliverable', ''),
                'status':      p.get('status', 'ready'),
                'brief':       p.get('brief', ''),
                'agents':      json.dumps(p.get('agents', [])),
                'created_at':  p.get('createdAt', now),
                'updated_at':  now,
            })
            con.commit()
            con.close()
            self._json_response({'ok': True})
        except Exception as e:
            self._json_response({'error': str(e)}, 500)

    # ── Deliverable Library API ───────────────────────────────────────────────

    def _get_deliverables(self):
        try:
            # Parse optional query params: ?client=X&verdict=Y&q=search
            parsed = urllib.parse.urlparse(self.path)
            params = urllib.parse.parse_qs(parsed.query)
            client  = params.get('client',  [None])[0]
            verdict = params.get('verdict', [None])[0]
            q       = params.get('q',       [None])[0]

            sql    = "SELECT * FROM deliverable_library WHERE 1=1"
            args   = []
            if client:
                sql += " AND client LIKE ?"
                args.append(f'%{client}%')
            if verdict:
                sql += " AND verdict = ?"
                args.append(verdict)
            if q:
                sql += " AND (label LIKE ? OR client LIKE ? OR content LIKE ?)"
                args += [f'%{q}%', f'%{q}%', f'%{q}%']
            sql += " ORDER BY created_at DESC"

            con  = db_connect()
            rows = con.execute(sql, args).fetchall()
            con.close()
            # Don't send full content in listing — return preview only
            result = []
            for r in rows:
                d = dict(r)
                d['preview'] = d['content'][:400] + ('…' if len(d['content']) > 400 else '')
                result.append(d)
            self._json_response(result)
        except Exception as e:
            self._json_response({'error': str(e)}, 500)

    def _save_deliverable(self):
        try:
            p   = self._read_json_body()
            now = datetime.datetime.utcnow().isoformat()
            content = p.get('content', '')
            word_count = len(content.split())
            con = db_connect()
            con.execute("""
                INSERT INTO deliverable_library
                    (id, engagement_id, gate_id, label, content, verdict,
                     client, sector, deliverable_type, word_count, created_at)
                VALUES
                    (:id,:engagement_id,:gate_id,:label,:content,:verdict,
                     :client,:sector,:deliverable_type,:word_count,:created_at)
                ON CONFLICT(id) DO UPDATE SET
                    content    = excluded.content,
                    verdict    = excluded.verdict,
                    word_count = excluded.word_count
            """, {
                'id':              p.get('id'),
                'engagement_id':   p.get('engagement_id', ''),
                'gate_id':         p.get('gate_id', ''),
                'label':           p.get('label', 'Deliverable'),
                'content':         content,
                'verdict':         p.get('verdict', ''),
                'client':          p.get('client', ''),
                'sector':          p.get('sector', ''),
                'deliverable_type':p.get('deliverable_type', ''),
                'word_count':      word_count,
                'created_at':      now,
            })
            con.commit()
            con.close()
            self._json_response({'ok': True, 'word_count': word_count})
        except Exception as e:
            self._json_response({'error': str(e)}, 500)

    def _delete_deliverable(self, del_id):
        try:
            con = db_connect()
            con.execute("DELETE FROM deliverable_library WHERE id = ?", (del_id,))
            con.commit()
            con.close()
            self._json_response({'ok': True})
        except Exception as e:
            self._json_response({'error': str(e)}, 500)

    # ── Anthropic proxy ───────────────────────────────────────────────────────

    def _proxy_anthropic(self):
        length = int(self.headers.get('Content-Length', 0))
        body   = self.rfile.read(length)

        fwd = {
            'Content-Type':      self.headers.get('Content-Type', 'application/json'),
            'anthropic-version': self.headers.get('anthropic-version', '2023-06-01'),
        }
        if self.headers.get('x-api-key'):
            fwd['x-api-key'] = self.headers['x-api-key']
        if self.headers.get('anthropic-beta'):
            fwd['anthropic-beta'] = self.headers['anthropic-beta']

        req = urllib.request.Request(ANTHROPIC_URL, data=body, headers=fwd, method='POST')
        try:
            with urllib.request.urlopen(req) as resp:
                data = resp.read()
                self.send_response(resp.status)
                self._cors_headers()
                self.send_header('Content-Type',   'application/json')
                self.send_header('Content-Length', str(len(data)))
                self.end_headers()
                self.wfile.write(data)
        except urllib.error.HTTPError as e:
            data = e.read()
            self.send_response(e.code)
            self._cors_headers()
            self.send_header('Content-Type',   'application/json')
            self.send_header('Content-Length', str(len(data)))
            self.end_headers()
            self.wfile.write(data)
        except Exception as e:
            msg = json.dumps({'error': {'message': str(e)}}).encode()
            self.send_response(502)
            self._cors_headers()
            self.send_header('Content-Type',   'application/json')
            self.send_header('Content-Length', str(len(msg)))
            self.end_headers()
            self.wfile.write(msg)

    # ── Document generation ───────────────────────────────────────────────────

    def _generate(self, fmt):
        try:
            length  = int(self.headers.get('Content-Length', 0))
            payload = json.loads(self.rfile.read(length))
            content = payload.get('content', '')
            title   = payload.get('title',   'APEX Deliverable')
            client  = payload.get('client',  '')
            sector  = payload.get('sector',  '')

            if fmt == 'docx':
                data     = _build_docx(content, title, client, sector)
                mime     = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
                filename = _slug(title) + '.docx'
            else:
                data     = _build_pptx(content, title, client, sector)
                mime     = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
                filename = _slug(title) + '.pptx'

            self.send_response(200)
            self._cors_headers()
            self.send_header('Content-Type',        mime)
            self.send_header('Content-Disposition', f'attachment; filename="{filename}"')
            self.send_header('Content-Length',      str(len(data)))
            self.end_headers()
            self.wfile.write(data)

        except Exception as e:
            import traceback; traceback.print_exc()
            msg = json.dumps({'error': str(e)}).encode()
            self.send_response(500)
            self._cors_headers()
            self.send_header('Content-Type',   'application/json')
            self.send_header('Content-Length', str(len(msg)))
            self.end_headers()
            self.wfile.write(msg)


if __name__ == '__main__':
    os.chdir(STATIC_DIR)
    init_db()
    server = http.server.HTTPServer(('', PORT), ProxyHandler)
    print(f'APEX Console  ->  http://localhost:{PORT}')
    print(f'Proxy target  ->  {ANTHROPIC_URL}')
    print('Endpoints: /api/messages  /api/generate-docx  /api/generate-pptx')
    print('           /api/engagements  /api/deliverables')
    print('Press Ctrl+C to stop.\n')
    try:
        server.serve_forever()
    except KeyboardInterrupt:
        print('\nStopped.')
        sys.exit(0)
