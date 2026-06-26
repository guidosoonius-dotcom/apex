"""
APEX API — single Flask entrypoint for Vercel's Python runtime.
Consolidates: engagements, deliverables, messages, generate-docx, generate-pptx.
Also serves the console SPA for all non-API routes.
"""
import sys, os, pathlib
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from flask import Flask, request, jsonify, Response, send_file
from _shared import db, parse_md, strip_inline, is_meta, slug
import json, io, datetime, urllib.parse, urllib.request, urllib.error

CONSOLE_DIR = pathlib.Path(__file__).parent.parent / 'console'

app = Flask(__name__)

ANTHROPIC_URL = 'https://api.anthropic.com/v1/messages'

CORS = {
    'Access-Control-Allow-Origin':  '*',
    'Access-Control-Allow-Methods': 'GET, POST, DELETE, OPTIONS',
    'Access-Control-Allow-Headers': 'Content-Type, x-api-key, anthropic-version, anthropic-beta',
}

NAVY  = (15,  27,  45)
TEAL  = (14, 124, 123)
GOLD  = (201, 168,  76)
WHITE = (240, 245, 250)
FOG   = (139, 163, 190)


@app.after_request
def add_cors(resp):
    for k, v in CORS.items():
        resp.headers[k] = v
    return resp


def options():
    return Response(status=204)


# ── Engagements ───────────────────────────────────────────────────────────────

@app.route('/api/engagements', methods=['OPTIONS'])
def engagements_options():
    return options()

@app.route('/api/engagements', methods=['GET'])
def get_engagements():
    try:
        con  = db()
        rows = con.execute('SELECT * FROM engagements ORDER BY created_at DESC').fetchall()
        con.close()
        return jsonify([dict(r) for r in rows])
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/engagements', methods=['POST'])
def post_engagement():
    try:
        p   = request.json
        now = datetime.datetime.utcnow().isoformat()
        con = db()
        con.execute("""
            INSERT INTO engagements
                (id, client, sector, deliverable, status, brief, agents, created_at, updated_at)
            VALUES
                (:id,:client,:sector,:deliverable,:status,:brief,:agents,:now,:now)
            ON CONFLICT(id) DO UPDATE SET
                status     = excluded.status,
                updated_at = excluded.updated_at
        """, {
            'id':          p.get('id'),
            'client':      p.get('client', ''),
            'sector':      p.get('sector', ''),
            'deliverable': p.get('deliverable', ''),
            'status':      p.get('status', 'ready'),
            'brief':       p.get('brief', ''),
            'agents':      json.dumps(p.get('agents', [])),
            'now':         now,
        })
        con.commit()
        con.close()
        return jsonify({'ok': True})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ── Deliverables ──────────────────────────────────────────────────────────────

@app.route('/api/deliverables', methods=['OPTIONS'])
def deliverables_options():
    return options()

@app.route('/api/deliverables', methods=['GET'])
def get_deliverables():
    try:
        client  = request.args.get('client')
        verdict = request.args.get('verdict')
        q       = request.args.get('q')
        del_id  = request.args.get('id')
        con     = db()

        if del_id:
            row = con.execute('SELECT * FROM deliverable_library WHERE id = ?', (del_id,)).fetchone()
            con.close()
            return (jsonify(dict(row)) if row else (jsonify({'error': 'Not found'}), 404))

        sql, args = 'SELECT * FROM deliverable_library WHERE 1=1', []
        if client:
            sql += ' AND client LIKE ?';  args.append(f'%{client}%')
        if verdict:
            sql += ' AND verdict = ?';    args.append(verdict)
        if q:
            sql += ' AND (label LIKE ? OR client LIKE ? OR content LIKE ?)';  args += [f'%{q}%']*3
        sql += ' ORDER BY created_at DESC'

        rows = con.execute(sql, args).fetchall()
        con.close()
        result = []
        for r in rows:
            d = dict(r)
            d['preview'] = d['content'][:400] + ('…' if len(d['content']) > 400 else '')
            result.append(d)
        return jsonify(result)
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/deliverables', methods=['POST'])
def post_deliverable():
    try:
        p       = request.json
        now     = datetime.datetime.utcnow().isoformat()
        content = p.get('content', '')
        wc      = len(content.split())
        con     = db()
        con.execute("""
            INSERT INTO deliverable_library
                (id, engagement_id, gate_id, label, content, verdict,
                 client, sector, deliverable_type, word_count, created_at)
            VALUES
                (:id,:engagement_id,:gate_id,:label,:content,:verdict,
                 :client,:sector,:deliverable_type,:word_count,:now)
            ON CONFLICT(id) DO UPDATE SET
                content    = excluded.content,
                verdict    = excluded.verdict,
                word_count = excluded.word_count
        """, {
            'id':              p.get('id'),
            'engagement_id':   p.get('engagement_id', ''),
            'gate_id':         p.get('gate_id', ''),
            'label':           p.get('label', 'Deliverable'),
            'content':         content,
            'verdict':         p.get('verdict', ''),
            'client':          p.get('client', ''),
            'sector':          p.get('sector', ''),
            'deliverable_type':p.get('deliverable_type', ''),
            'word_count':      wc,
            'now':             now,
        })
        con.commit()
        con.close()
        return jsonify({'ok': True, 'word_count': wc})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/deliverables/<path:del_id>', methods=['OPTIONS'])
def deliverable_id_options(del_id):
    return options()

@app.route('/api/deliverables/<path:del_id>', methods=['GET'])
def get_deliverable(del_id):
    del_id = urllib.parse.unquote(del_id)
    try:
        con = db()
        row = con.execute('SELECT * FROM deliverable_library WHERE id = ?', (del_id,)).fetchone()
        con.close()
        return (jsonify(dict(row)) if row else (jsonify({'error': 'Not found'}), 404))
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/deliverables/<path:del_id>', methods=['DELETE'])
def delete_deliverable(del_id):
    del_id = urllib.parse.unquote(del_id)
    try:
        con = db()
        con.execute('DELETE FROM deliverable_library WHERE id = ?', (del_id,))
        con.commit()
        con.close()
        return jsonify({'ok': True})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ── Messages — Anthropic proxy ────────────────────────────────────────────────

@app.route('/api/messages', methods=['OPTIONS'])
def messages_options():
    return options()

@app.route('/api/messages', methods=['POST'])
def post_messages():
    body = request.get_data()
    fwd  = {'Content-Type': request.headers.get('Content-Type', 'application/json'),
             'anthropic-version': request.headers.get('anthropic-version', '2023-06-01')}
    if request.headers.get('x-api-key'):
        fwd['x-api-key'] = request.headers['x-api-key']
    if request.headers.get('anthropic-beta'):
        fwd['anthropic-beta'] = request.headers['anthropic-beta']

    req = urllib.request.Request(ANTHROPIC_URL, data=body, headers=fwd, method='POST')
    try:
        with urllib.request.urlopen(req) as resp:
            return Response(resp.read(), status=resp.status, content_type='application/json')
    except urllib.error.HTTPError as e:
        return Response(e.read(), status=e.code, content_type='application/json')
    except Exception as e:
        return jsonify({'error': str(e)}), 502


# ── Generate DOCX ─────────────────────────────────────────────────────────────

def build_docx(content, title, client, sector):
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor
    import re

    doc = Document()
    for sec in doc.sections:
        sec.top_margin    = Inches(1.0)
        sec.bottom_margin = Inches(1.0)
        sec.left_margin   = Inches(1.25)
        sec.right_margin  = Inches(1.25)

    def rgb(*t): return RGBColor(*t)

    def run(para, text, bold=False, color=None, size=None, italic=False):
        r = para.add_run(text)
        r.bold   = bold
        r.italic = italic
        if color: r.font.color.rgb = rgb(*color)
        if size:  r.font.size = Pt(size)
        return r

    def inline(para, text, size=11, default_color=None):
        parts = re.split(r'\*\*(.+?)\*\*', text)
        for i, part in enumerate(parts):
            if part:
                run(para, part, bold=(i % 2 == 1), size=size, color=default_color)

    cover = doc.add_paragraph()
    cover.alignment = 1
    run(cover, title, bold=True, color=TEAL, size=28)
    sub = doc.add_paragraph()
    sub.alignment = 1
    run(sub, '  ·  '.join(filter(None, [client, sector])), color=FOG, size=12)
    doc.add_paragraph()

    for btype, btext in parse_md(content):
        if btype == 'title':
            pass
        elif btype == 'h1':
            h = doc.add_heading(level=1); h.clear()
            run(h, strip_inline(btext), bold=True, color=TEAL, size=16)
        elif btype == 'h2':
            h = doc.add_heading(level=2); h.clear()
            run(h, strip_inline(btext), bold=True, color=GOLD, size=13)
        elif btype == 'bullet':
            p = doc.add_paragraph(style='List Bullet')
            inline(p, strip_inline(btext), size=11)
        elif btype == 'hr':
            doc.add_paragraph()
        elif btype == 'para':
            if is_meta(btext): continue
            p = doc.add_paragraph()
            inline(p, strip_inline(btext), size=11)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.read()


@app.route('/api/generate-docx', methods=['OPTIONS'])
def generate_docx_options():
    return options()

@app.route('/api/generate-docx', methods=['POST'])
def generate_docx():
    try:
        p        = request.json
        content  = p.get('content', '')
        title    = p.get('title',   'APEX Deliverable')
        client   = p.get('client',  '')
        sector   = p.get('sector',  '')
        data     = build_docx(content, title, client, sector)
        mime     = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        resp     = Response(data, status=200, content_type=mime)
        resp.headers['Content-Disposition'] = f'attachment; filename="{slug(title)}.docx"'
        return resp
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ── Generate PPTX ─────────────────────────────────────────────────────────────

def build_pptx(content, title, client, sector):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def rgb(*t): return RGBColor(*t)

    def bg(slide, color=NAVY):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = rgb(*color)

    def accent_bar(slide):
        s = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.07), Inches(7.5))
        s.fill.solid(); s.fill.fore_color.rgb = rgb(*TEAL); s.line.fill.background()

    def add_textbox(slide, l, t, w, h):
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tb.text_frame.word_wrap = True
        return tb.text_frame

    def para(tf, text, size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT, bullet=False):
        p = tf.add_paragraph(); p.alignment = align
        if bullet: p.level = 1
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = rgb(*color)
        return p

    sl = prs.slides.add_slide(blank); bg(sl); accent_bar(sl)
    tf = add_textbox(sl, 0.5, 2.2, 12.5, 1.6)
    para(tf, title, size=38, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
    tf2 = add_textbox(sl, 0.5, 3.9, 12.5, 0.7)
    para(tf2, '  ·  '.join(filter(None, [client, sector])), size=16, color=FOG)
    tf3 = add_textbox(sl, 0.5, 6.8, 4, 0.4)
    para(tf3, 'APEX  ·  AI-native strategy consulting', size=10, color=FOG)

    sections, cur_h, cur_body = [], None, []
    for btype, btext in parse_md(content):
        if btype == 'h1':
            if cur_h is not None: sections.append((cur_h, cur_body))
            cur_h, cur_body = btext, []
        elif btype in ('title', 'hr', 'empty'): continue
        elif cur_h is not None: cur_body.append((btype, btext))
    if cur_h is not None: sections.append((cur_h, cur_body))

    for heading, body_lines in sections:
        sl = prs.slides.add_slide(blank); bg(sl); accent_bar(sl)
        tf_h = add_textbox(sl, 0.35, 0.25, 12.7, 0.85)
        para(tf_h, strip_inline(heading), size=22, bold=True, color=TEAL)
        bar = sl.shapes.add_shape(1, Inches(0.35), Inches(1.05), Inches(12.7), Pt(2))
        bar.fill.solid(); bar.fill.fore_color.rgb = rgb(*GOLD); bar.line.fill.background()
        tf_body = add_textbox(sl, 0.35, 1.25, 12.7, 5.9)
        chars = 0
        for btype, btext in body_lines:
            if chars > 1400: break
            clean = strip_inline(btext).strip()
            if not clean or is_meta(clean): continue
            if   btype == 'bullet': para(tf_body, '• ' + clean, size=14, color=WHITE, bullet=True)
            elif btype == 'h2':     para(tf_body, clean, size=14, bold=True, color=GOLD)
            elif btype == 'para':   para(tf_body, clean, size=13, color=WHITE)
            chars += len(clean)

    sl = prs.slides.add_slide(blank); bg(sl); accent_bar(sl)
    tf_e = add_textbox(sl, 0.5, 3.0, 12.5, 1.2)
    para(tf_e, 'APEX', size=48, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    tf_s = add_textbox(sl, 0.5, 4.1, 12.5, 0.6)
    para(tf_s, 'AI-native strategy consulting', size=16, color=FOG, align=PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


@app.route('/api/generate-pptx', methods=['OPTIONS'])
def generate_pptx_options():
    return options()

@app.route('/api/generate-pptx', methods=['POST'])
def generate_pptx():
    try:
        p       = request.json
        content = p.get('content', '')
        title   = p.get('title',   'APEX Deliverable')
        client  = p.get('client',  '')
        sector  = p.get('sector',  '')
        data    = build_pptx(content, title, client, sector)
        mime    = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        resp    = Response(data, status=200, content_type=mime)
        resp.headers['Content-Disposition'] = f'attachment; filename="{slug(title)}.pptx"'
        return resp
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ── SPA catch-all — serve console/index.html for all non-API routes ──────────

@app.route('/', defaults={'path': ''})
@app.route('/<path:path>')
def spa(path):
    index = CONSOLE_DIR / 'index.html'
    if index.exists():
        return send_file(str(index), mimetype='text/html')
    return 'Not found', 404
