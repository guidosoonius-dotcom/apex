"""
/api/generate-pptx — Convert markdown deliverable to branded PowerPoint deck.
"""
import sys, os
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from http.server import BaseHTTPRequestHandler
from _shared import CORS_HEADERS, send_cors_ok, read_json, parse_md, strip_inline, is_meta, slug
import json
import io

NAVY  = (15,  27,  45)
TEAL  = (14, 124, 123)
GOLD  = (201, 168,  76)
WHITE = (240, 245, 250)
FOG   = (139, 163, 190)


def build_pptx(content, title, client, sector):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def rgb(*t): return RGBColor(*t)

    def bg(slide, color=NAVY):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = rgb(*color)

    def accent_bar(slide):
        s = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.07), Inches(7.5))
        s.fill.solid(); s.fill.fore_color.rgb = rgb(*TEAL); s.line.fill.background()

    def add_textbox(slide, l, t, w, h):
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tb.text_frame.word_wrap = True
        return tb.text_frame

    def para(tf, text, size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT, bullet=False):
        p = tf.add_paragraph(); p.alignment = align
        if bullet: p.level = 1
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = rgb(*color)
        return p

    # Cover slide
    sl = prs.slides.add_slide(blank); bg(sl); accent_bar(sl)
    tf = add_textbox(sl, 0.5, 2.2, 12.5, 1.6)
    para(tf, title, size=38, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
    tf2 = add_textbox(sl, 0.5, 3.9, 12.5, 0.7)
    para(tf2, '  ·  '.join(filter(None, [client, sector])), size=16, color=FOG)
    tf3 = add_textbox(sl, 0.5, 6.8, 4, 0.4)
    para(tf3, 'APEX  ·  AI-native strategy consulting', size=10, color=FOG)

    # Content slides
    sections = []
    cur_h, cur_body = None, []
    for btype, btext in parse_md(content):
        if btype == 'h1':
            if cur_h is not None: sections.append((cur_h, cur_body))
            cur_h, cur_body = btext, []
        elif btype in ('title', 'hr', 'empty'): continue
        elif cur_h is not None: cur_body.append((btype, btext))
    if cur_h is not None: sections.append((cur_h, cur_body))

    for heading, body_lines in sections:
        sl = prs.slides.add_slide(blank); bg(sl); accent_bar(sl)
        tf_h = add_textbox(sl, 0.35, 0.25, 12.7, 0.85)
        para(tf_h, strip_inline(heading), size=22, bold=True, color=TEAL)
        bar = sl.shapes.add_shape(1, Inches(0.35), Inches(1.05), Inches(12.7), Pt(2))
        bar.fill.solid(); bar.fill.fore_color.rgb = rgb(*GOLD); bar.line.fill.background()
        tf_body = add_textbox(sl, 0.35, 1.25, 12.7, 5.9)
        chars = 0
        for btype, btext in body_lines:
            if chars > 1400: break
            clean = strip_inline(btext).strip()
            if not clean or is_meta(clean): continue
            if   btype == 'bullet': para(tf_body, '• ' + clean, size=14, color=WHITE, bullet=True)
            elif btype == 'h2':     para(tf_body, clean, size=14, bold=True, color=GOLD)
            elif btype == 'para':   para(tf_body, clean, size=13, color=WHITE)
            chars += len(clean)

    # End slide
    sl = prs.slides.add_slide(blank); bg(sl); accent_bar(sl)
    tf_e = add_textbox(sl, 0.5, 3.0, 12.5, 1.2)
    para(tf_e, 'APEX', size=48, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    tf_s = add_textbox(sl, 0.5, 4.1, 12.5, 0.6)
    para(tf_s, 'AI-native strategy consulting', size=16, color=FOG, align=PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


class handler(BaseHTTPRequestHandler):

    def log_message(self, *a): pass

    def do_OPTIONS(self):
        send_cors_ok(self)

    def do_POST(self):
        try:
            p       = read_json(self)
            content = p.get('content', '')
            title   = p.get('title',   'APEX Deliverable')
            client  = p.get('client',  '')
            sector  = p.get('sector',  '')
            data     = build_pptx(content, title, client, sector)
            mime     = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            filename = slug(title) + '.pptx'
            self.send_response(200)
            for k, v in CORS_HEADERS.items():
                self.send_header(k, v)
            self.send_header('Content-Type',        mime)
            self.send_header('Content-Disposition', f'attachment; filename="{filename}"')
            self.send_header('Content-Length',      str(len(data)))
            self.end_headers()
            self.wfile.write(data)
        except Exception as e:
            import traceback; traceback.print_exc()
            msg = json.dumps({'error': str(e)}).encode()
            self.send_response(500)
            for k, v in CORS_HEADERS.items():
                self.send_header(k, v)
            self.send_header('Content-Type',   'application/json')
            self.send_header('Content-Length', str(len(msg)))
            self.end_headers()
            self.wfile.write(msg)
